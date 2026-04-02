#!/usr/bin/env python3
"""
Generate cuRAMSES presentation (.pptx) from paper content.
Background: pre-dawn deep sky blue (#0c1628)
Focus: figures and concise text
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color palette: pre-dawn twilight sky ─────────────────────────
BG_DEEP    = RGBColor(0x0c, 0x16, 0x28)   # Deep sky background
BG_PANEL   = RGBColor(0x12, 0x1f, 0x35)   # Slightly lighter panel
ACCENT     = RGBColor(0xe8, 0x96, 0x5a)   # Dawn amber
ACCENT2    = RGBColor(0x6b, 0xa3, 0xd6)   # Pale sky blue
TEXT_MAIN  = RGBColor(0xec, 0xed, 0xf0)   # Off-white
TEXT_DIM   = RGBColor(0x8a, 0x95, 0xaa)   # Muted blue-gray
TEXT_EMPH  = RGBColor(0xff, 0xd7, 0x92)   # Warm highlight
WHITE      = RGBColor(0xff, 0xff, 0xff)
HORIZON    = RGBColor(0x2a, 0x1a, 0x3a)   # Faint purple horizon hint

# ── Slide dimensions: 16:9 ──────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MISC = '/gpfs/kjhan/Hydro/CUBE_HR5/code_cube/misc'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ── Helper functions ────────────────────────────────────────────

def set_bg(slide, color=BG_DEEP):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text,
             font_size=18, color=TEXT_MAIN, bold=False,
             alignment=PP_ALIGN.LEFT, font_name='Arial',
             line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf

def add_multiline(slide, left, top, width, height, lines,
                  font_size=18, color=TEXT_MAIN, bold=False,
                  alignment=PP_ALIGN.LEFT, line_spacing=1.4,
                  font_name='Arial'):
    """Add text box with multiple paragraphs from list of (text, kwargs_override) or str"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, kw = item, {}
        else:
            txt, kw = item[0], item[1]
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(kw.get('size', font_size))
        p.font.color.rgb = kw.get('color', color)
        p.font.bold = kw.get('bold', bold)
        p.font.name = kw.get('font', font_name)
        p.alignment = kw.get('align', alignment)
        ls = kw.get('line_spacing', line_spacing)
        p.line_spacing = Pt(kw.get('size', font_size) * ls)
        p.space_before = Pt(kw.get('space_before', 0))
        p.space_after = Pt(kw.get('space_after', 4))
    return tf

def add_image(slide, path, left, top, width=None, height=None):
    # Prefer PNG over PDF (pptx can't embed PDF)
    if path.endswith('.pdf'):
        png_path = path.replace('.pdf', '.png')
        if os.path.exists(png_path):
            path = png_path
    if not os.path.exists(path):
        add_text(slide, left, top, 4, 0.5, f'[Image: {os.path.basename(path)}]',
                 font_size=14, color=TEXT_DIM)
        return None
    kwargs = {}
    if width: kwargs['width'] = Inches(width)
    if height: kwargs['height'] = Inches(height)
    return slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)

def add_thin_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def title_slide(title, subtitle='', section_num=''):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    if section_num:
        add_text(slide, 0.8, 1.8, 11, 0.6, section_num,
                 font_size=16, color=ACCENT2, font_name='Arial')
    add_text(slide, 0.8, 2.3, 11, 2.0, title,
             font_size=42, color=TEXT_MAIN, bold=True, font_name='Georgia')
    add_thin_line(slide, 0.8, 4.5, 3.5, ACCENT)
    if subtitle:
        add_text(slide, 0.8, 4.8, 11, 1.5, subtitle,
                 font_size=20, color=TEXT_DIM, font_name='Arial')
    return slide

def content_slide(title, section_label=''):
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide)
    # Title bar
    add_text(slide, 0.8, 0.3, 10, 0.8, title,
             font_size=30, color=TEXT_MAIN, bold=True, font_name='Georgia')
    add_thin_line(slide, 0.8, 1.1, 2.5, ACCENT)
    if section_label:
        add_text(slide, 10.5, 0.4, 2.5, 0.4, section_label,
                 font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.RIGHT)
    return slide

def figure_slide(title, img_path, caption='', section_label='',
                 img_left=1.5, img_top=1.5, img_width=10, img_height=None):
    slide = content_slide(title, section_label)
    add_image(slide, img_path, img_left, img_top, img_width, img_height)
    if caption:
        add_text(slide, 0.8, 6.5, 11.5, 0.8, caption,
                 font_size=14, color=TEXT_DIM)
    return slide


# ════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, 0.8, 1.5, 11, 0.6, 'cuRAMSES',
         font_size=56, color=ACCENT, bold=True, font_name='Georgia')
add_text(s, 0.8, 2.5, 11, 1.5,
         'Scalable AMR Optimizations\nfor Large-Scale Cosmological Simulations',
         font_size=34, color=TEXT_MAIN, font_name='Georgia', line_spacing=1.3)
add_thin_line(s, 0.8, 4.3, 5, ACCENT)
add_text(s, 0.8, 4.7, 11, 0.6, 'Juhan Kim',
         font_size=22, color=TEXT_EMPH, font_name='Arial')
add_text(s, 0.8, 5.3, 11, 0.8,
         'Center for Advanced Computation\nKorea Institute for Advanced Study',
         font_size=16, color=TEXT_DIM, font_name='Arial')

# ════════════════════════════════════════════════════════════════
# SLIDE 2: Outline
# ════════════════════════════════════════════════════════════════
s = content_slide('Outline')
items = [
    '1.  Motivation: Scaling Challenges in AMR Cosmology',
    '2.  K-Section Domain Decomposition',
    '3.  Hierarchical MPI Exchange',
    '4.  Morton Key Hash Table & Memory Savings',
    '5.  Poisson Solver: Multigrid + FFTW3',
    '6.  Feedback Spatial Binning',
    '7.  Variable-NCPU Restart',
    '8.  Performance Benchmarks',
    '9.  GPU Acceleration',
    '10. Conclusions & Outlook',
]
add_multiline(s, 1.2, 1.5, 10, 5.5, items, font_size=22, color=TEXT_MAIN,
              line_spacing=1.8)

# ════════════════════════════════════════════════════════════════
# SLIDES 3-5: Motivation
# ════════════════════════════════════════════════════════════════
title_slide('Motivation',
            'Scaling cosmological AMR to the exascale regime',
            'SECTION 1')

s = content_slide('The Challenge', '1. Motivation')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Modern cosmological simulations demand:', {'size': 22, 'color': TEXT_EMPH}),
    ('',{'size': 8}),
    ('\u2022  Large volume (> 1 Gpc/h)\u00b3  for statistical precision', {}),
    ('\u2022  Deep AMR hierarchy (dynamic range > 10\u2075)', {}),
    ('\u2022  Sub-parsec resolution for SMBH feedback, ISM, star formation', {}),
    ('',{'size': 12}),
    ('Bottlenecks in standard RAMSES:', {'size': 22, 'color': ACCENT, 'space_before': 12}),
    ('',{'size': 8}),
    ('\u2022  MPI_ALLTOALL  \u2192  O(N\u00b2_rank) communication', {}),
    ('\u2022  Hilbert curve: any rank may border any other rank', {}),
    ('\u2022  Memory: nbor array > 1 GB for N_gridmax = 5M', {}),
    ('\u2022  Multigrid Poisson solver \u2248 50% of wallclock', {}),
    ('\u2022  One-file-per-rank I/O prevents flexible restarts', {}),
], font_size=20, color=TEXT_MAIN, line_spacing=1.5)

s = content_slide('cuRAMSES: A Suite of Solutions', '1. Motivation')
add_multiline(s, 1.0, 1.5, 5.5, 5.5, [
    ('K-section decomposition', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  Neighbour-only P2P exchange', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('Morton key hash table', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  O(1) lookup, >190 MB/rank saved', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('Multigrid + FFTW3', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  Merged sweeps, 8.3\u00d7 base speedup', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('Feedback spatial binning', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  260\u00d7 SNII, 30\u00d7 AGN speedup', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('Variable-N_rank restart', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  HDF5 + binary, arbitrary rank count', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('GPU offloading', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('  MPI+OpenMP+CUDA, auto dispatch', {'size': 16, 'color': TEXT_DIM}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)
add_image(s, f'{MISC}/domain_decomposition_improved.png', 7.0, 1.5, width=5.8)

# ════════════════════════════════════════════════════════════════
# SLIDES 6-9: K-Section Domain Decomposition
# ════════════════════════════════════════════════════════════════
title_slide('K-Section Domain Decomposition',
            'Recursive hierarchical spatial partitioning\nreplacing Hilbert curve ordering',
            'SECTION 2')

figure_slide('Progressive K-Section Decomposition',
             f'{MISC}/ksection_progressive.png',
             'N_rank = 12 = 3 \u00d7 2 \u00d7 2: split into 3 slabs, bisect each, bisect again. '
             'Tree structure encodes both domain hierarchy and communication pattern.',
             '2. K-Section',
             img_left=1.0, img_top=1.4, img_width=11)

s = content_slide('Memory-Weighted Load Balancing', '2. K-Section')
add_multiline(s, 1.0, 1.5, 5.5, 5.0, [
    ('Cost function per cell:', {'size': 20, 'color': ACCENT}),
    ('', {'size': 6}),
    ('C = (w_grid + n_part \u00b7 w_part + n_sink \u00b7 w_sink) / 8',
     {'size': 20, 'color': TEXT_EMPH, 'font': 'Courier New'}),
    ('', {'size': 12}),
    ('Default RAMSES:', {'size': 18, 'color': TEXT_DIM}),
    ('  cost = 80 + N_part   (underweights particles 3\u20134\u00d7)', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 8}),
    ('cuRAMSES:', {'size': 18, 'color': TEXT_MAIN}),
    ('  w_grid = 2256 B (nvar=14)', {'size': 16}),
    ('  w_part = 12 B/particle', {'size': 16}),
    ('  w_sink = 500 (computational weight)', {'size': 16}),
    ('', {'size': 12}),
    ('Result:', {'size': 20, 'color': ACCENT}),
    ('  Peak/mean memory ratio: 2.5 \u2192 1.3', {'size': 18, 'color': TEXT_EMPH}),
    ('  M_max/M_min \u2264 1.05 across 2\u201364 ranks', {'size': 18, 'color': TEXT_EMPH}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)
add_image(s, f'{MISC}/fig_memory_balance.pdf', 7.2, 1.5, width=5.5)

# ════════════════════════════════════════════════════════════════
# SLIDES 10-12: Hierarchical Exchange
# ════════════════════════════════════════════════════════════════
title_slide('Auto-Tuning MPI Communications',
            'Hierarchical exchange eliminates global collectives',
            'SECTION 3')

figure_slide('Hierarchical K-Section Exchange',
             f'{MISC}/hierarchical_exchange.png',
             'N_rank=12: each rank communicates with at most \u03a3(k_l \u2212 1) = 4 partners, '
             'not N_rank\u22121 = 11. Communication count is O(\u03a3 k_l), independent of total rank count.',
             '3. MPI Exchange',
             img_left=1.5, img_top=1.4, img_width=10)

s = content_slide('Communication Complexity', '3. MPI Exchange')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Three exchange backends with auto-tuning:', {'size': 22, 'color': ACCENT}),
    ('', {'size': 10}),
    ('MPI_ALLTOALLV', {'size': 20, 'bold': True}),
    ('  \u2022  Count: O(N_rank)     Buffer: O(N_rank)', {'size': 16, 'color': TEXT_DIM, 'font': 'Courier New'}),
    ('  \u2022  Simple, but all-to-all metadata exchange', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 6}),
    ('Point-to-Point (P2P)', {'size': 20, 'bold': True}),
    ('  \u2022  Count: O(N_nb)       Buffer: O(N_nb)', {'size': 16, 'color': TEXT_DIM, 'font': 'Courier New'}),
    ('  \u2022  Neighbour-only, but needs discovery phase', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 6}),
    ('K-Section Hierarchical', {'size': 20, 'bold': True, 'color': TEXT_EMPH}),
    ('  \u2022  Count: O(\u03a3 k_l)    Buffer: O(k_max \u00b7 N_gh)', {'size': 16, 'color': TEXT_DIM, 'font': 'Courier New'}),
    ('  \u2022  Tree-guided, constant partner count', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 10}),
    ('Auto-tune: each component independently selects the fastest backend',
     {'size': 18, 'color': ACCENT2}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDES 13-15: Morton Key Hash Table
# ════════════════════════════════════════════════════════════════
title_slide('Morton Key Hash Table',
            'O(1) neighbour lookup replacing the nbor array\nwith >190 MB/rank memory savings',
            'SECTION 4')

s = content_slide('The nbor Problem & Morton Solution', '4. Morton Key')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Problem: nbor array', {'size': 22, 'color': ACCENT}),
    ('  \u2022  6 neighbour pointers per grid \u00d7 N_gridmax', {'size': 18}),
    ('  \u2022  >1 GB for N_gridmax = 5M', {'size': 18}),
    ('  \u2022  Complex maintenance during refinement/migration', {'size': 18}),
    ('', {'size': 12}),
    ('Solution: Morton key hash table', {'size': 22, 'color': ACCENT}),
    ('  \u2022  Encode grid position as Morton (Z-order) key', {'size': 18}),
    ('  \u2022  Neighbour = bit manipulation on the key', {'size': 18}),
    ('  \u2022  Per-level hash table with O(1) lookup', {'size': 18}),
    ('  \u2022  No nbor array needed at all', {'size': 18}),
    ('', {'size': 12}),
    ('Savings:', {'size': 22, 'color': TEXT_EMPH}),
    ('  >190 MB/rank  +  simpler code  +  faster refinement', {'size': 20, 'color': TEXT_EMPH}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDES 16-19: Poisson Solver
# ════════════════════════════════════════════════════════════════
title_slide('Poisson Solver Optimizations',
            'Multigrid V-cycle + FFTW3 direct solver',
            'SECTION 5')

figure_slide('Optimized V-Cycle Algorithm',
             f'{MISC}/fig_mg_vcycle.pdf',
             'Merged red-black sweeps and fused residual-norm reduce ghost-zone exchanges from 9 to 5 per level. '
             'Poisson share reduced from 50% to 40% of total runtime.',
             '5. Poisson',
             img_left=2.0, img_top=1.4, img_width=9)

s = content_slide('FFTW3 Direct Poisson Solver', '5. Poisson')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Base-level Poisson: multigrid \u2192 FFTW3 FFT', {'size': 22, 'color': ACCENT}),
    ('', {'size': 8}),
    ('Two modes:', {'size': 20}),
    ('  \u2022  N \u2264 256\u00b3: ALLREDUCE + local FFT', {'size': 18}),
    ('  \u2022  N > 256\u00b3: Sparse P2P + FFTW3 MPI slab decomposition', {'size': 18}),
    ('', {'size': 8}),
    ('Sparse P2P exchange:', {'size': 20, 'color': ACCENT2}),
    ('  \u2022  Only 5\u20138 partner ranks (not N_rank)', {'size': 18}),
    ('  \u2022  bisec_cpubox partner cache for O(1) routing', {'size': 18}),
    ('  \u2022  In-place R2C saves ~90 MB/rank', {'size': 18}),
    ('', {'size': 12}),
    ('Performance:', {'size': 22, 'color': TEXT_EMPH}),
    ('  Base-level: 241s \u2192 28.9s  (8.3\u00d7 speedup)', {'size': 22, 'color': TEXT_EMPH}),
    ('  No GPU required for this acceleration', {'size': 18, 'color': TEXT_DIM}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDES 20-22: Feedback Spatial Binning
# ════════════════════════════════════════════════════════════════
title_slide('Feedback Spatial Binning',
            'Hash-based spatial indexing for SN and AGN feedback',
            'SECTION 6')

figure_slide('Spatial Hash Binning',
             f'{MISC}/fig_spatial_binning.pdf',
             'Domain partitioned into uniform bin grid. Each target cell checks 27 neighbouring bins (3D). '
             'Complexity: O(N_cells \u00d7 27 \u00b7 n_SN/bin) instead of O(N_cells \u00d7 N_SN).',
             '6. Feedback',
             img_left=2.5, img_top=1.4, img_width=8)

s = content_slide('Feedback Performance', '6. Feedback')
add_multiline(s, 1.0, 1.5, 5.5, 5.0, [
    ('Speedup results:', {'size': 22, 'color': ACCENT}),
    ('', {'size': 10}),
    ('SNII feedback:', {'size': 20, 'bold': True}),
    ('  260\u00d7 speedup', {'size': 28, 'color': TEXT_EMPH, 'bold': True}),
    ('', {'size': 8}),
    ('AGN feedback:', {'size': 20, 'bold': True}),
    ('  30\u00d7 speedup', {'size': 28, 'color': TEXT_EMPH, 'bold': True}),
    ('', {'size': 12}),
    ('Sink merger: Oct-tree FoF', {'size': 20, 'color': ACCENT2}),
    ('  O(N log N) vs O(N\u00b2) brute-force', {'size': 16, 'color': TEXT_DIM}),
    ('  600\u00d7 at N_sink = 10\u2075', {'size': 18, 'color': TEXT_EMPH}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)
add_image(s, f'{MISC}/fig_fof_scaling.png', 7.0, 1.5, width=5.5)

# ════════════════════════════════════════════════════════════════
# SLIDES 23-24: Variable-NCPU Restart
# ════════════════════════════════════════════════════════════════
title_slide('Variable-N_CPU Restart',
            'HDF5 parallel I/O with arbitrary rank count',
            'SECTION 7')

s = content_slide('HDF5 Schema & Restart Flexibility', '7. Variable-NCPU')
add_multiline(s, 1.0, 1.5, 5.0, 5.0, [
    ('Single shared HDF5 file:', {'size': 20, 'color': ACCENT}),
    ('  \u2022  All ranks write via MPI-IO', {'size': 16}),
    ('  \u2022  AMR/hydro/gravity per level', {'size': 16}),
    ('  \u2022  Particles as flat arrays', {'size': 16}),
    ('', {'size': 10}),
    ('Restart compatibility:', {'size': 20, 'color': ACCENT}),
    ('  \u2022  Binary \u2194 HDF5', {'size': 16}),
    ('  \u2022  Hilbert \u2194 K-section', {'size': 16}),
    ('  \u2022  Any N_rank \u2192 any N_rank', {'size': 16}),
    ('', {'size': 10}),
    ('All 8 format combinations supported:', {'size': 18, 'color': TEXT_EMPH}),
    ('  {Binary, HDF5} \u00d7 {Hilbert, K-sec}', {'size': 18, 'color': TEXT_EMPH}),
    ('  \u00d7 {Hilbert, K-sec} (file \u00d7 run)', {'size': 18, 'color': TEXT_EMPH}),
], font_size=16, color=TEXT_MAIN, line_spacing=1.4)
add_image(s, f'{MISC}/fig_hdf5_schema.pdf', 6.5, 1.3, width=6.2)

# ════════════════════════════════════════════════════════════════
# SLIDES 25-31: Performance Benchmarks
# ════════════════════════════════════════════════════════════════
title_slide('Performance Benchmarks',
            'Strong scaling, OpenMP, memory balance,\nmulti-node, and weak scaling',
            'SECTION 8')

figure_slide('Strong Scaling',
             f'{MISC}/fig_strong_scaling.png',
             'Top: Cosmo512 single node, 33.9\u00d7 at 64 ranks (53% efficiency). '
             'Bottom: Cosmo1024 multi-node, 26.5\u00d7 at 32 nodes (83% efficiency).',
             '8. Performance',
             img_left=1.5, img_top=1.3, img_width=10)

figure_slide('OpenMP Thread Scaling',
             f'{MISC}/omp_scaling.png',
             'N_rank=4 on dual-socket AMD EPYC 7543 (64 cores). '
             '5.8\u00d7 overall, 10.5\u00d7 for MG solver. OMP parallelism is highly effective.',
             '8. Performance',
             img_left=2.0, img_top=1.3, img_width=9)

figure_slide('Hybrid MPI/OpenMP Scaling (8 nodes)',
             f'{MISC}/fig_cosmo1024_omp_scaling.png',
             '512 cores on 8 nodes. Super-linear speedup from rank reduction at 2\u20138 threads. '
             'Production config: 64 ranks \u00d7 8 threads.',
             '8. Performance',
             img_left=2.0, img_top=1.3, img_width=9)

figure_slide('Weak Scaling',
             f'{MISC}/fig_weak_scaling.png',
             'Grammar cluster, nthread=8. Efficiency: 88.5% at 128 cores, 65.7% at 1024 cores. '
             'Dominated by MPI collective overhead at high rank count.',
             '8. Performance',
             img_left=2.0, img_top=1.3, img_width=9)

s = content_slide('Performance Summary', '8. Performance')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Single 64-core node (Cosmo512):', {'size': 22, 'color': ACCENT}),
    ('  33.9\u00d7 speedup  (53% parallel efficiency)', {'size': 22, 'color': TEXT_EMPH}),
    ('', {'size': 10}),
    ('OpenMP (16 threads/rank):', {'size': 22, 'color': ACCENT}),
    ('  5.8\u00d7 overall,  10.5\u00d7 MG solver', {'size': 22, 'color': TEXT_EMPH}),
    ('', {'size': 10}),
    ('Multi-node Cosmo1024 (3.6\u00d710\u2078 grids, 1.15\u00d710\u2079 particles):', {'size': 22, 'color': ACCENT}),
    ('  1\u201332 nodes: 26.5\u00d7 speedup (83% efficiency)', {'size': 22, 'color': TEXT_EMPH}),
    ('  Near-ideal scaling to 16 nodes (98%)', {'size': 22, 'color': TEXT_EMPH}),
    ('', {'size': 10}),
    ('Memory balance: M_max/M_min \u2264 1.05 across all configurations',
     {'size': 20, 'color': ACCENT2}),
    ('All diagnostics match Hilbert reference within roundoff',
     {'size': 18, 'color': TEXT_DIM}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDES 32-36: GPU Offloading
# ════════════════════════════════════════════════════════════════
title_slide('GPU Acceleration',
            'MPI + OpenMP + CUDA hybrid dispatch\nwith GPU-resident mesh data',
            'SECTION 9')

s = content_slide('GPU Architecture', '9. GPU')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Three-level parallelism:', {'size': 22, 'color': ACCENT}),
    ('  MPI (inter-node)  +  OpenMP (intra-node)  +  CUDA (GPU)', {'size': 18}),
    ('', {'size': 10}),
    ('Dynamic dispatch model:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  Auto-tune selects CPU or GPU per component', {'size': 18}),
    ('  \u2022  GPU fallback to CPU if memory insufficient', {'size': 18}),
    ('', {'size': 8}),
    ('GPU-resident mesh data:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  uold, unew, rho, phi, f persist on GPU', {'size': 18}),
    ('  \u2022  Only ghost zones transferred per step', {'size': 18}),
    ('  \u2022  H2D overhead: 51% \u2192 5.7%', {'size': 18, 'color': TEXT_EMPH}),
    ('', {'size': 8}),
    ('GPU multigrid:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  Gauss-Seidel red-black kernels', {'size': 18}),
    ('  \u2022  Restriction/Prolongation on GPU', {'size': 18}),
    ('  \u2022  cuFFT for base-level direct solve', {'size': 18}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

figure_slide('GPU Performance Model',
             f'{MISC}/fig_gpu_mg_model.png',
             '(a) MG time vs CPU-GPU bandwidth. (b) Speedup shows 1.95\u00d7 asymptotic limit. '
             '(c) Speedup vs r (OMP threads per GPU). H100 NVL: 1.70\u00d7 MG, net 1.21\u00d7 overall.',
             '9. GPU',
             img_left=1.0, img_top=1.3, img_width=11)

s = content_slide('GPU Assessment', '9. GPU')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('Current results (PCIe-connected GPUs):', {'size': 22, 'color': ACCENT}),
    ('', {'size': 6}),
    ('  MG Poisson:  1.70\u00d7 (H100 NVL),  2.06\u00d7 (A100 SXM4)', {'size': 20, 'color': TEXT_EMPH}),
    ('  Godunov:     0.69\u00d7 (slower \u2014 PCIe bandwidth limited)', {'size': 20, 'color': RGBColor(0xff, 0x88, 0x88)}),
    ('  Net overall: 1.21\u00d7 improvement', {'size': 20, 'color': TEXT_EMPH}),
    ('', {'size': 12}),
    ('Why limited?', {'size': 20, 'color': ACCENT2}),
    ('  \u2022  AMR stencil: irregular octree access \u2192 latency-bound', {'size': 18}),
    ('  \u2022  Godunov: 23 GB cudaMemcpy per level per step', {'size': 18}),
    ('  \u2022  SM-count limited, not HBM bandwidth limited', {'size': 18}),
    ('', {'size': 12}),
    ('Future: NVLink-C2C (GH200) predicted ~2\u00d7 MG, >5\u00d7 with GPU-direct halo exchange',
     {'size': 18, 'color': ACCENT2}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDES 37-39: Conclusions
# ════════════════════════════════════════════════════════════════
title_slide('Conclusions',
            'Summary and future directions',
            'SECTION 10')

s = content_slide('Key Contributions', '10. Conclusions')
add_multiline(s, 1.0, 1.5, 11, 5.5, [
    ('1. K-section decomposition', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   O(\u03a3 k_l) comm, memory balance <5% imbalance', {'size': 16, 'color': TEXT_DIM}),
    ('2. Morton key hash table', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   >190 MB/rank saved, O(1) neighbour lookup', {'size': 16, 'color': TEXT_DIM}),
    ('3. Multigrid optimization', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   44% fewer exchanges, Poisson 50% \u2192 40%', {'size': 16, 'color': TEXT_DIM}),
    ('4. FFTW3 direct solver', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   8.3\u00d7 base-level speedup, sparse P2P', {'size': 16, 'color': TEXT_DIM}),
    ('5. Feedback spatial binning', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   260\u00d7 SNII, 30\u00d7 AGN, 600\u00d7 FoF merger', {'size': 16, 'color': TEXT_DIM}),
    ('6. Variable-N_rank restart', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   HDF5/binary \u00d7 Hilbert/K-section, any N_rank', {'size': 16, 'color': TEXT_DIM}),
    ('7. GPU acceleration', {'size': 20, 'color': ACCENT, 'bold': True}),
    ('   1.70\u00d7 MG (H100), 1.21\u00d7 net; GH200 pathway to 2\u00d7+', {'size': 16, 'color': TEXT_DIM}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.15)

s = content_slide('Conservation & Verification', '10. Conclusions')
add_multiline(s, 1.0, 1.5, 11, 5.0, [
    ('All optimizations preserve physics:', {'size': 24, 'color': ACCENT}),
    ('', {'size': 10}),
    ('\u2022  e_cons, e_pot, e_kin, e_int match Hilbert reference within roundoff', {'size': 20}),
    ('\u2022  m_tot, n_star, n_sink identical', {'size': 20}),
    ('\u2022  Deterministic cell-index-seeded RNG', {'size': 20}),
    ('   (star formation independent of DD and thread schedule)', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 12}),
    ('Exception: merged red-black sweep \u2192 0.5% e_cons shift', {'size': 18, 'color': ACCENT2}),
    ('(intentional: stale boundary values, well-understood)', {'size': 16, 'color': TEXT_DIM}),
    ('', {'size': 16}),
    ('Scaling headline numbers:', {'size': 22, 'color': ACCENT}),
    ('  Single node:   33.9\u00d7  (64 cores, 53% eff.)', {'size': 20, 'color': TEXT_EMPH}),
    ('  Multi-node:    26.5\u00d7  (32 nodes, 83% eff.)', {'size': 20, 'color': TEXT_EMPH}),
    ('  Weak scaling:  65.7%  at 1024 cores', {'size': 20, 'color': TEXT_EMPH}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

s = content_slide('Future Directions', '10. Conclusions')
add_multiline(s, 1.0, 1.5, 11, 5.0, [
    ('Scaling verification:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  Test at \u226510\u2074 MPI ranks', {'size': 18}),
    ('  \u2022  K-section designed for O(10\u2075) concurrency', {'size': 18}),
    ('', {'size': 10}),
    ('GPU pathway:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  NVLink-C2C (GH200): ~2\u00d7 MG predicted', {'size': 18}),
    ('  \u2022  GPU-direct halo exchange: >5\u00d7 potential', {'size': 18}),
    ('  \u2022  Device-resident restriction/prolongation', {'size': 18}),
    ('', {'size': 10}),
    ('Non-standard cosmology:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  Massive neutrinos via Fourier transfer functions', {'size': 18}),
    ('  \u2022  Coupled dark energy perturbations', {'size': 18}),
    ('', {'size': 10}),
    ('Production:', {'size': 22, 'color': ACCENT}),
    ('  \u2022  cuRAMSES in use for next-generation simulation project', {'size': 18}),
    ('  \u2022  Public release upon benchmark completion', {'size': 18}),
], font_size=18, color=TEXT_MAIN, line_spacing=1.3)

# ════════════════════════════════════════════════════════════════
# SLIDE 40: Thank you
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s)
add_text(s, 0.8, 2.0, 11, 1.5, 'Thank You',
         font_size=52, color=TEXT_EMPH, bold=True, font_name='Georgia',
         alignment=PP_ALIGN.CENTER)
add_thin_line(s, 4.5, 3.5, 4, ACCENT)
add_text(s, 0.8, 4.0, 11, 0.8, 'kjhan@kias.re.kr',
         font_size=22, color=ACCENT2, font_name='Arial',
         alignment=PP_ALIGN.CENTER)
add_text(s, 0.8, 5.0, 11, 1.0,
         'Center for Advanced Computation\nKorea Institute for Advanced Study',
         font_size=18, color=TEXT_DIM, font_name='Arial',
         alignment=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────
outpath = f'{MISC}/cuRAMSES_presentation.pptx'
prs.save(outpath)
print(f'Saved: {outpath}')
print(f'Total slides: {len(prs.slides)}')
